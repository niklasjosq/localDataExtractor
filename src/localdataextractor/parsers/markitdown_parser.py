from __future__ import annotations

from pathlib import Path

from localdataextractor.models import ParsedResult, SourceReference, TableBlock
from localdataextractor.parsers.base import ParserContext
from localdataextractor.parsers.common import markdown_to_blocks_and_tables


class MarkItDownParser:
    name = "markitdown"

    def extract(self, path: Path, context: ParserContext) -> ParsedResult:
        warnings: list[str] = []
        markdown_text = ""

        try:
            from markitdown import MarkItDown  # type: ignore

            converter = MarkItDown()
            result = converter.convert(str(path))
            markdown_text = getattr(result, "text_content", "") or str(result)
        except Exception as exc:  # pragma: no cover - optional dependency path
            warnings.append(f"MarkItDown unavailable/failed: {exc}")

        if not markdown_text.strip():
            fallback_text, fallback_tables, fallback_warnings = _fallback_extract(path)
            warnings.extend(fallback_warnings)
            blocks, tables = markdown_to_blocks_and_tables(fallback_text)
            if fallback_tables:
                tables.extend(fallback_tables)
            return ParsedResult(
                parser_name=self.name,
                warnings=warnings,
                blocks=blocks,
                tables=tables,
            )

        blocks, tables = markdown_to_blocks_and_tables(markdown_text)
        return ParsedResult(
            parser_name=self.name,
            warnings=warnings,
            blocks=blocks,
            tables=tables,
        )


def _fallback_extract(path: Path) -> tuple[str, list[TableBlock], list[str]]:
    ext = path.suffix.lower()
    warnings: list[str] = []

    if ext in {".xlsx", ".xls"}:
        text, tables, warns = _fallback_spreadsheet(path)
        return text, tables, warns
    if ext == ".docx":
        text, tables, warns = _fallback_docx(path)
        return text, tables, warns
    if ext == ".pptx":
        text, warns = _fallback_pptx(path)
        return text, [], warns

    data = path.read_text(encoding="utf-8", errors="replace")
    return data, [], warnings


def _fallback_docx(path: Path) -> tuple[str, list[TableBlock], list[str]]:
    warnings: list[str] = []
    tables: list[TableBlock] = []
    lines: list[str] = []
    try:
        from docx import Document  # type: ignore

        doc = Document(str(path))
        for para in doc.paragraphs:
            text = para.text.strip()
            if text:
                lines.append(text)
        for idx, table in enumerate(doc.tables, start=1):
            headers: list[str] = []
            body: list[list[str]] = []
            for row_index, row in enumerate(table.rows):
                values = [cell.text.strip() for cell in row.cells]
                if row_index == 0:
                    headers = values
                else:
                    body.append(values)
            tables.append(
                TableBlock(
                    table_id=f"tbl_{idx:05d}",
                    header_rows=[headers] if headers else [],
                    body_rows=body,
                    column_count=max([len(headers)] + [len(r) for r in body]),
                    row_count=len(body),
                    confidence=88.0,
                )
            )
            if headers:
                lines.append("| " + " | ".join(headers) + " |")
                lines.append("| " + " | ".join(["---"] * len(headers)) + " |")
                for row in body:
                    lines.append("| " + " | ".join(row) + " |")
    except Exception as exc:  # pragma: no cover
        warnings.append(f"python-docx fallback failed: {exc}")
    return "\n\n".join(lines), tables, warnings


def _fallback_pptx(path: Path) -> tuple[str, list[str]]:
    warnings: list[str] = []
    lines: list[str] = []
    try:
        from pptx import Presentation  # type: ignore

        prs = Presentation(str(path))
        for slide_num, slide in enumerate(prs.slides, start=1):
            lines.append(f"## Slide {slide_num}")
            for shape in slide.shapes:
                text = getattr(shape, "text", "").strip()
                if text:
                    lines.append(text)
    except Exception as exc:  # pragma: no cover
        warnings.append(f"python-pptx fallback failed: {exc}")
    return "\n\n".join(lines), warnings


def _fallback_spreadsheet(path: Path) -> tuple[str, list[TableBlock], list[str]]:
    ext = path.suffix.lower()
    if ext == ".xlsx":
        return _fallback_xlsx(path)
    return _fallback_xls(path)


def _fallback_xlsx(path: Path) -> tuple[str, list[TableBlock], list[str]]:
    warnings: list[str] = []
    tables: list[TableBlock] = []
    lines: list[str] = []
    try:
        from openpyxl import load_workbook  # type: ignore

        wb = load_workbook(filename=path, read_only=True, data_only=True)
        table_id = 0
        for sheet in wb.worksheets:
            lines.append(f"## Sheet: {sheet.title}")
            rows: list[list[str]] = []
            for row in sheet.iter_rows(values_only=True):
                values = ["" if v is None else str(v) for v in row]
                if any(cell.strip() for cell in values):
                    rows.append(values)
            if not rows:
                continue
            headers = rows[0]
            body = rows[1:]
            table_id += 1
            tables.append(
                TableBlock(
                    table_id=f"tbl_{table_id:05d}",
                    source=SourceReference(sheet=sheet.title),
                    header_rows=[headers],
                    body_rows=body,
                    column_count=max(len(r) for r in rows),
                    row_count=len(body),
                    confidence=90.0,
                )
            )
            lines.append("| " + " | ".join(headers) + " |")
            lines.append("| " + " | ".join(["---"] * len(headers)) + " |")
            for row in body:
                lines.append("| " + " | ".join(row) + " |")
    except Exception as exc:  # pragma: no cover
        warnings.append(f"openpyxl fallback failed: {exc}")
    return "\n".join(lines), tables, warnings


def _fallback_xls(path: Path) -> tuple[str, list[TableBlock], list[str]]:
    warnings: list[str] = []
    tables: list[TableBlock] = []
    lines: list[str] = []
    try:
        import xlrd  # type: ignore

        wb = xlrd.open_workbook(str(path))
        table_id = 0
        for sheet in wb.sheets():
            lines.append(f"## Sheet: {sheet.name}")
            rows: list[list[str]] = []
            for i in range(sheet.nrows):
                row = [str(sheet.cell_value(i, j)) for j in range(sheet.ncols)]
                if any(cell.strip() for cell in row):
                    rows.append(row)
            if not rows:
                continue
            headers = rows[0]
            body = rows[1:]
            table_id += 1
            tables.append(
                TableBlock(
                    table_id=f"tbl_{table_id:05d}",
                    source=SourceReference(sheet=sheet.name),
                    header_rows=[headers],
                    body_rows=body,
                    column_count=sheet.ncols,
                    row_count=len(body),
                    confidence=86.0,
                )
            )
            lines.append("| " + " | ".join(headers) + " |")
            lines.append("| " + " | ".join(["---"] * len(headers)) + " |")
            for row in body:
                lines.append("| " + " | ".join(row) + " |")
    except Exception as exc:  # pragma: no cover
        warnings.append(f"xlrd fallback failed: {exc}")
    return "\n".join(lines), tables, warnings
